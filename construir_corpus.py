# -*- coding: utf-8 -*-
"""
Constructor del corpus canonico - Historia Critica del Peru (122005-A, 2026-1)

Fase 1 del metodo: ingesta.
Recorre carpeta_documentos, extrae texto anclado a la fuente (archivo+pagina),
resuelve duplicados pptx/pdf, incorpora OCR de escaneados, y produce:

  corpus/corpus.jsonl   una linea por bloque, con procedencia completa
  corpus/manifiesto.json  inventario de documentos con metadatos
  corpus/reporte.txt      resumen legible

No indexa ni recupera todavia. Solo construye la representacion persistente.
"""
import sys, io, os, json, re, hashlib, glob, unicodedata
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace", line_buffering=True)

import fitz
import tiktoken
from pptx import Presentation
import docx

RAIZ = os.path.dirname(os.path.abspath(__file__))
DOCS = os.path.join(RAIZ, "carpeta_documentos", "Historia crítica del Perú")
OUT = os.path.join(RAIZ, "corpus")
OCR_DIR = os.path.join(OUT, "_ocr")
os.makedirs(OUT, exist_ok=True)

ENC = tiktoken.get_encoding("cl100k_base")
ntok = lambda s: len(ENC.encode(s)) if s else 0

CURSO = {
    "codigo": "122005",
    "nombre": "Historia Crítica del Perú",
    "seccion": "A",
    "periodo": "2026-1",
    "profesor": "Juan Fonseca",
    "departamento": "Humanidades",
    "creditos": 4,
}

# Estructura oficial extraida del cronograma (fuente: Cronograma_HCP_2026_1_A.pdf)
SEMANAS = {
    1:  {"unidad": 1, "fechas": "16-22 marzo", "tema": "Conceptos demográficos y principales indicadores"},
    2:  {"unidad": 1, "fechas": "23-29 marzo", "tema": "Transiciones demográficas y epidemiológicas"},
    3:  {"unidad": 1, "fechas": "30 marzo-5 abril", "tema": "Geografía del Perú: Andes, cuencas, El Niño"},
    4:  {"unidad": 1, "fechas": "6-12 abril", "tema": "Organización del territorio; ciudades y redes urbanas"},
    5:  {"unidad": 1, "fechas": "13-19 abril", "tema": "Esquema urbano-rural y redes urbanas"},
    6:  {"unidad": 1, "fechas": "20-26 abril", "tema": "Cambio climático, desastres y políticas de desarrollo"},
    7:  {"unidad": 1, "fechas": "27 abril-3 mayo", "tema": "El caso de la región Áncash"},
    8:  {"unidad": None, "fechas": "4-10 mayo", "tema": "EXÁMENES PARCIALES (5 de mayo)"},
    9:  {"unidad": 2, "fechas": "11-17 mayo", "tema": "La idea de nación; construcción del Estado nacional"},
    10: {"unidad": 2, "fechas": "18-24 mayo", "tema": "Regiones y macrorregiones; la costa norte colonial"},
    11: {"unidad": 2, "fechas": "25-31 mayo", "tema": "Independencia: norte patriota y sur realista"},
    12: {"unidad": 2, "fechas": "1-7 junio", "tema": "Haciendas azucareras y APRA; sur andino prehispánico"},
    13: {"unidad": 2, "fechas": "8-14 junio", "tema": "Sur andino colonial; Confederación Perú-Boliviana"},
    14: {"unidad": 2, "fechas": "15-21 junio", "tema": "Arequipa s.XIX; movimientos indígenas s.XX"},
}
UNIDADES = {
    1: "Demografía, territorio y medioambiente en el Perú",
    2: "Procesos socioeconómicos y recursos naturales en los espacios regionales",
}

# Bibliografia oficial: mapea archivo -> (cita, caracter). Fuente: cronograma.
BIBLIO = {
 "Aramburú_2012_Poblacióndiferente_pags-55-72": ("Aramburú, C. E. (2012). Una población diferente: cinco décadas de cambio demográfico. En Cambios sociales en el Perú 1968-2008, pp. 55-72.", "obligatoria"),
 "Amat_y_León_2012_Perúnuestrodecadadía": ("Amat y León, C. (2012). Somos un país de montañas tropicales. En El Perú nuestro de cada día, cap. 1, pp. 14-35.", "obligatoria"),
 "moneda-190-07": ("Guabloche, J. y Gutiérrez, A. P. (2020). Crecimiento de la población peruana y estructura demográfica. Revista Moneda 190, pp. 36-40. BCRP.", "obligatoria"),
 "Contreras_2020_Crisisdemográfica_sigloXVI": ("Contreras, C. (2020). La crisis demográfica del siglo XVI en los Andes. Diálogo Andino 61, pp. 7-25.", "obligatoria"),
 "Contreras_1994_Orígenes_explosión_demográfica": ("Contreras, C. (1994). Sobre los orígenes de la explosión demográfica en el Perú: 1876-1940. Documento de Trabajo 61, IEP.", "obligatoria"),
 "Aramburú_Mendoza_2015_Futuro_población_peruana": ("Aramburú, C. y Mendoza, W. (2015). El futuro de la población peruana. Debates en Sociología 41, pp. 5-24.", "obligatoria"),
 "Torero_Escobal_2000_Como_enfrentar_una_geografia_adversa": ("Torero, M. y Escobal, J. (2000). Geografía y desarrollo. En ¿Cómo enfrentar una geografía adversa?, pp. 6-13.", "obligatoria"),
 "Cavieses_2001_Imprints_El_Niño": ("Caviedes, C. N. (2001). Imprints of El Niño in World History. En El Niño in History, pp. 172-215.", "obligatoria"),
 "Marzal_Ludeña_2017_Ciudadesintermedias": ("Marzal, V. y Ludeña, W. (2017). Perú: ciudades intermedias en el sistema urbano regional y nacional, pp. 220-257.", "obligatoria"),
 "Sobrevilla_2021_Disputa_jurisdicciones": ("Sobrevilla Perea, N. (2021). Disputa por las jurisdicciones y la formación estatal, 1808-1850, pp. 19-58.", "obligatoria"),
 "Contreras_2002_Centralismo_peruano": ("Contreras, C. (2002). El centralismo peruano en su perspectiva histórica. IEP.", "recomendada"),
 "Espinoza_Fort_Espinoza_2022_Reorganizar_el Perú": ("Espinoza, A., Fort, R. y Espinoza, M. (2022). Reorganizar el Perú: ciudades intermedias y desarrollo. GRADE.", "obligatoria"),
 "Carey_2014_Glaciares_CambioClimático_caps_3_4": ("Carey, M. (2014). Glaciares, cambio climático y desastres naturales. IEP, caps. 3-4.", "obligatoria"),
 "Contreras_sf_Concepto_Nación": ("Contreras, C. (s.f.). El concepto de nación y su evolución a través del tiempo. BCRP.", "obligatoria"),
 "Contreras_Cueto_2008_Caminos_Ciencia_Estado": ("Contreras, C. y Cueto, M. (2008). Caminos, ciencia y Estado en el Perú, 1850-1930. HCS-Manguinhos 15(3), pp. 635-655.", "obligatoria"),
 "Aldana_Pereyra_2022_Regiones vivas y activas": ("Aldana, S. y Pereyra, N. (2022). Regiones vivas y activas. Proyecto Especial Bicentenario.", "obligatoria"),
 "RivasplataVarillas_2014_Cambio_paisajes_costa_norte": ("Rivasplata Varillas, P. (2014). Cambio de paisajes de la costa norte peruana. Historia 2.0, IV(8), pp. 47-73.", "obligatoria"),
 "O_Phelan_2019_Norte_patriota_Sur_realista": ("O'Phelan Godoy, S. (2019). El norte patriota y el sur realista: la división territorial del Perú (1820-1824).", "obligatoria"),
 "Klarén_1976_Haciendas_Azucareras_Apra": ("Klarén, P. (1976). Formación de las haciendas azucareras y orígenes del APRA. IEP, cap. 3.", "obligatoria"),
 "Cahill_1993_Independencia_sociedad_fiscalidad": ("Cahill, D. (1993). Independencia, sociedad y fiscalidad. El Sur Andino (1780-1880). RCHA 19.", "obligatoria"),
 "Renique_2022_Nación_Radical_cap.4": ("Rénique, J. L. (2022). La nación radical, cap. 4: Indios e indigenistas en el Altiplano sur andino. La Siniestra.", "obligatoria"),
}

def norm(s):
    """Normaliza texto: colapsa espacios, une guionizacion de fin de linea."""
    if not s: return ""
    s = s.replace("­", "")
    s = re.sub(r"(\w)-\n(\w)", r"\1\2", s)      # palabra cortada por guion
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

# Artefactos de digitalizacion. Los PDFs escaneados del corpus vienen de
# CamScanner y arrastran su marca en cada pagina; sin limpiarla aparece
# dentro de las citas. Las lineas sin palabras reales son ruido de OCR
# (bordes de pagina, sombras del escaneo) y solo gastan presupuesto.
MARCAS_OCR = re.compile(r"(?im)^.*(escaneado con camscanner|scanned with camscanner|camscanner).*$")

def limpiar_ocr(t):
    """Quita marcas de digitalizacion y lineas que son ruido de escaneo.

    El criterio es la proporcion de caracteres que forman parte de palabras
    reales. Una linea como '1 1 1i 1 1 on r 1 Artiatin niii 1i' contiene una
    secuencia larga de letras, pero casi todo su contenido son fragmentos
    sueltos: sin este ratio pasaria el filtro y gastaria presupuesto.
    """
    t = MARCAS_OCR.sub("", t)
    salida = []
    for ln in t.split("\n"):
        s = ln.strip()
        if len(s) > 20:
            en_palabras = sum(len(w) for w in re.findall(r"[a-zA-ZáéíóúñÁÉÍÓÚÑ]{3,}", s))
            if en_palabras / len(s) < 0.5:
                continue      # menos de la mitad del texto son palabras: ruido
        salida.append(ln)
    return norm("\n".join(salida))

def sha(p):
    h = hashlib.sha256()
    with open(p, "rb") as f:
        for c in iter(lambda: f.read(1 << 20), b""): h.update(c)
    return h.hexdigest()[:16]

def clasificar(rel, stem):
    """Deriva tipo, unidad, semana y autoridad desde la ruta."""
    partes = rel.split(os.sep)
    unidad = semana = None
    for p in partes:
        m = re.match(r"Unidad (\d)", p)
        if m: unidad = int(m.group(1))
        m = re.match(r"Semana (\d+)", p)
        if m: semana = int(m.group(1))
    if "Trabajos calificados" in rel:
        tipo, autoridad = "evaluacion", "oficial"
    elif stem.startswith("Clase") or "Clase_" in stem:
        tipo, autoridad = "clase", "docente"
    elif stem in ("Cronograma_HCP_2026_1_A",) or "122005" in stem:
        tipo, autoridad = "rector", "oficial"
    elif stem.startswith("HCP_"):
        tipo, autoridad = "guion", "oficial"
    else:
        tipo, autoridad = "lectura", "academica"
    return tipo, unidad, semana, autoridad

def autores_de(stem, tipo):
    if tipo != "lectura": return []
    cab = re.split(r"_\d{4}|_sf_", stem)[0]
    cab = cab.replace("_", " ").strip()
    return [a for a in re.split(r"\s+", cab) if len(a) > 2][:3]

# Anios que el nombre de archivo no declara pero el cronograma si.
# moneda-190-07 es Guabloche y Gutierrez 2020; Contreras "s.f." no tiene
# fecha declarada en la fuente y se deja como None a proposito.
ANIO_MANUAL = {"moneda-190-07": 2020}

def anio_de(stem):
    if stem in ANIO_MANUAL: return ANIO_MANUAL[stem]
    m = re.search(r"_(\d{4})_", stem + "_") or re.search(r"(\d{4})", stem)
    if "_sf_" in stem: return None
    if m:
        y = int(m.group(1))
        if 1900 < y < 2030: return y
    return None

# ---------- extractores: devuelven [(unidad_n, texto, extra)] ----------
def ext_pdf(p):
    d = fitz.open(p); out = []
    for i, pg in enumerate(d):
        t = norm(pg.get_text("text"))
        out.append((i + 1, t, {"imagenes": len(pg.get_images(full=True))}))
    d.close(); return out

def ext_pptx(p):
    prs = Presentation(p); out = []
    for i, s in enumerate(prs.slides):
        partes, imgs, tablas = [], 0, 0
        for sh in s.shapes:
            if sh.shape_type is not None and "PICTURE" in str(sh.shape_type): imgs += 1
            if getattr(sh, "has_table", False):
                tablas += 1
                for r in sh.table.rows:
                    fila = [c.text.strip() for c in r.cells if c.text.strip()]
                    if fila: partes.append(" | ".join(fila))
            if sh.has_text_frame and sh.text_frame.text.strip():
                partes.append(sh.text_frame.text.strip())
        notas = ""
        if s.has_notes_slide:
            notas = norm(s.notes_slide.notes_text_frame.text)
        out.append((i + 1, norm("\n".join(partes)), {"imagenes": imgs, "tablas": tablas, "notas": notas}))
    return out

def ext_docx(p):
    d = docx.Document(p); partes = []
    for pa in d.paragraphs:
        if pa.text.strip(): partes.append(pa.text.strip())
    for t in d.tables:
        for r in t.rows:
            fila = [c.text.strip() for c in r.cells if c.text.strip()]
            if fila: partes.append(" | ".join(fila))
    return [(1, norm("\n".join(partes)), {})]

def ext_ocr(stem):
    f = os.path.join(OCR_DIR, stem + ".json")
    if not os.path.exists(f): return None
    data = json.load(open(f, encoding="utf-8"))
    return [(pg["page"], limpiar_ocr(pg["text"]), {"ocr": True}) for pg in data["paginas"]]

# ---------- inventario y resolucion de duplicados ----------
archivos = []
for p in glob.glob(os.path.join(DOCS, "**", "*"), recursive=True):
    if not os.path.isfile(p): continue
    ext = os.path.splitext(p)[1].lower()
    if ext not in (".pdf", ".pptx", ".docx"): continue
    archivos.append(p)

por_stem = {}
for p in archivos:
    rel = os.path.relpath(p, DOCS)
    stem = os.path.splitext(rel)[0]
    por_stem.setdefault(stem, {})[os.path.splitext(p)[1].lower()] = p

print(f"Archivos detectados: {len(archivos)}  |  documentos logicos: {len(por_stem)}\n")

# Decision de canonicidad (medida en el sondeo: el PDF conserva mas texto que el PPTX)
def elegir_canonico(v):
    if ".pdf" in v and ".pptx" in v: return ".pdf", [".pptx"]
    if ".pdf" in v and ".docx" in v: return ".pdf", [".docx"]
    e = list(v.keys())[0]; return e, []

manifiesto, bloques = [], []
bid = 0
for stem in sorted(por_stem):
    v = por_stem[stem]
    ext_can, alt = elegir_canonico(v)
    p = v[ext_can]
    rel = os.path.relpath(p, DOCS)
    base = os.path.basename(stem)
    tipo, unidad, semana, autoridad = clasificar(rel, base)

    if ext_can == ".pdf":   unidades = ext_pdf(p);  ulabel = "pagina"
    elif ext_can == ".pptx": unidades = ext_pptx(p); ulabel = "diapositiva"
    else:                    unidades = ext_docx(p); ulabel = "documento"

    # Si existe OCR para este documento, comparamos con la extraccion nativa y
    # nos quedamos con la que mas texto aporte. No basta con exigir que el
    # nativo sea casi vacio: hay PDFs mixtos, con texto en algunas paginas y
    # solo imagen en el resto (Sobrevilla tenia 10.436 tokens nativos frente a
    # 38.416 por OCR, con la misma calidad de palabras).
    chars = sum(len(t) for _, t, _ in unidades)
    ocr_usado = False
    o = ext_ocr(base)
    if o:
        chars_ocr = sum(len(t) for _, t, _ in o)
        if chars_ocr > chars * 1.15:      # margen para no cambiar por ruido
            unidades, ocr_usado, ulabel = o, True, "pagina"
            chars = chars_ocr

    cita, caracter = BIBLIO.get(base, (None, None))
    doc_id = re.sub(r"[^a-zA-Z0-9]+", "-", unicodedata.normalize("NFKD", base)
                    .encode("ascii", "ignore").decode()).strip("-").lower()[:60]

    toks = 0
    for n, t, extra in unidades:
        if not t.strip(): continue
        tk = ntok(t)
        toks += tk
        bid += 1
        bloques.append({
            "bloque_id": f"{doc_id}#{ulabel[0]}{n}",
            "doc_id": doc_id,
            "texto": t,
            "tokens": tk,
            # --- anclaje a la fuente (procedencia) ---
            "fuente": {"archivo": rel, "formato": ext_can, ulabel: n,
                       "ocr": bool(extra.get("ocr", False))},
            # --- metadatos estructurales, derivados de la ruta ---
            "curso": CURSO["codigo"], "unidad": unidad, "semana": semana,
            "tipo": tipo, "autoridad": autoridad,
            "imagenes": extra.get("imagenes", 0),
            "notas_presentador": extra.get("notas", ""),
        })

    manifiesto.append({
        "doc_id": doc_id, "titulo": base, "archivo": rel,
        "formato_canonico": ext_can, "formatos_descartados": alt,
        "tipo": tipo, "unidad": unidad,
        "unidad_nombre": UNIDADES.get(unidad), "semana": semana,
        "semana_tema": SEMANAS.get(semana, {}).get("tema"),
        "semana_fechas": SEMANAS.get(semana, {}).get("fechas"),
        "autoridad": autoridad, "caracter": caracter,
        "autores": autores_de(base, tipo), "anio": anio_de(base),
        "cita": cita, "unidades_fisicas": len(unidades), "etiqueta_unidad": ulabel,
        "tokens": toks, "ocr": ocr_usado, "hash": sha(p),
        "mb": round(os.path.getsize(p) / 1e6, 2),
    })
    est = "OCR" if ocr_usado else "  "
    print(f"[{est}] {toks:>7} tok  U{unidad or '-'} S{semana or '-'}  {tipo:<10} {base[:56]}")

# --- escritura ---
with open(os.path.join(OUT, "corpus.jsonl"), "w", encoding="utf-8") as f:
    for b in bloques: f.write(json.dumps(b, ensure_ascii=False) + "\n")
with open(os.path.join(OUT, "manifiesto.json"), "w", encoding="utf-8") as f:
    json.dump({"curso": CURSO, "unidades": UNIDADES, "semanas": SEMANAS,
               "documentos": manifiesto}, f, ensure_ascii=False, indent=2)

# --- reporte ---
tot = sum(m["tokens"] for m in manifiesto)
L = []
L.append("CORPUS - Historia Critica del Peru (122005-A, 2026-1) - prof. Juan Fonseca")
L.append("=" * 78)
L.append(f"Documentos logicos : {len(manifiesto)}   (de {len(archivos)} archivos)")
L.append(f"Bloques anclados   : {len(bloques)}")
L.append(f"Tokens totales     : {tot:,}")
L.append(f"Documentos con OCR : {sum(1 for m in manifiesto if m['ocr'])}")
L.append("")
L.append("POR TIPO")
for t in sorted({m["tipo"] for m in manifiesto}):
    g = [m for m in manifiesto if m["tipo"] == t]
    L.append(f"  {t:<12} {len(g):>3} docs  {sum(x['tokens'] for x in g):>9,} tok")
L.append("")
L.append("POR SEMANA")
for s in sorted(SEMANAS):
    g = [m for m in manifiesto if m["semana"] == s]
    tk = sum(x["tokens"] for x in g)
    marca = "" if g else "   <-- SIN MATERIAL"
    L.append(f"  S{s:<2} U{SEMANAS[s]['unidad'] or '-'}  {len(g):>2} docs {tk:>8,} tok  {SEMANAS[s]['tema'][:44]}{marca}")
sin_sem = [m for m in manifiesto if m["semana"] is None]
L.append(f"  (sin semana: {len(sin_sem)} docs, {sum(x['tokens'] for x in sin_sem):,} tok - rectores y evaluaciones)")
L.append("")
L.append("DOCUMENTOS MAS GRANDES")
for m in sorted(manifiesto, key=lambda x: -x["tokens"])[:8]:
    L.append(f"  {m['tokens']:>8,} tok  {m['titulo'][:60]}")
rep = "\n".join(L)
open(os.path.join(OUT, "reporte.txt"), "w", encoding="utf-8").write(rep)
print("\n" + rep)
print(f"\n-> corpus/corpus.jsonl  ({len(bloques)} bloques)")
print(f"-> corpus/manifiesto.json")
print(f"-> corpus/reporte.txt")
