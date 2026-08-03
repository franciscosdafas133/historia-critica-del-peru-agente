# -*- coding: utf-8 -*-
"""
Fase A: extractor con pistas fisicas.

A diferencia de construir_corpus.py (produccion, usa page.get_text("text"),
solo texto plano SIN bbox/fuente/indentacion -- confirmado por exploracion),
este extractor usa page.get_text("dict") para PDF, y lee posicion/fuente
directamente de python-pptx/python-docx, para poder inducir estructura
organizacional a partir de pistas orto-tipograficas reales.

No importa nada de construir_corpus.py (aislamiento total, ver plan) --
la logica de recorrer carpetas se reimplementa aqui, aceptando algo de
duplicacion a cambio de independencia completa del pipeline de produccion.
"""
import os
import re

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Emu
import docx

from forma_ir.tipos import Bloque, mayusculas_ratio, termina_en_puntuacion

# Bit 4 (valor 16) de span["flags"] en PyMuPDF = negrita.
# Confirmado empiricamente sobre Cronograma_HCP_2026_1_A.pdf: todo texto en
# fuente "Arial,Bold" trae flags=16; texto en "Calibri" (no negrita) trae
# flags=0. Bit 1 (valor 2) es italica en la convencion de PyMuPDF, pero se
# verifica por separado via el nombre de fuente tambien (ver _es_italica).
_BIT_NEGRITA = 1 << 4


def _es_negrita(flags: int, nombre_fuente: str) -> bool:
    if flags & _BIT_NEGRITA:
        return True
    return "bold" in nombre_fuente.lower()


def _es_italica(flags: int, nombre_fuente: str) -> bool:
    if flags & (1 << 1):
        return True
    nf = nombre_fuente.lower()
    return "italic" in nf or "oblique" in nf


def _mismo_parrafo(linea_a: dict, linea_b: dict, interlineado_max_pt: float = 4.0) -> bool:
    """True si `linea_b` es la continuacion visual de `linea_a` dentro del
    MISMO parrafo -- misma fuente/tamano/negrita/cursiva (misma firma
    tipografica) y salto vertical entre ellas compatible con interlineado
    normal (no un salto de parrafo).

    Se exige ademas indentacion.x0 compatible (delta pequeno) para no
    fusionar lineas de columnas DISTINTAS que por coincidencia comparten
    tamano de fuente en un layout multi-columna -- esto era precisamente
    la causa raiz encontrada en Checkpoint C3a: sin este chequeo, lineas
    de columna izquierda y columna derecha en un PDF de dos columnas se
    fusionarian igual si su y0 casi coincidiera."""
    sa, sb = linea_a["spans"][0], linea_b["spans"][0]
    if round(sa["size"], 2) != round(sb["size"], 2):
        return False
    if _es_negrita(sa["flags"], sa["font"]) != _es_negrita(sb["flags"], sb["font"]):
        return False
    if _es_italica(sa["flags"], sa["font"]) != _es_italica(sb["flags"], sb["font"]):
        return False
    bbox_a, bbox_b = linea_a["bbox"], linea_b["bbox"]
    if abs(bbox_a[0] - bbox_b[0]) > 3.0:
        return False  # x0 distinto -> probablemente otra columna, no continuacion
    salto = bbox_b[1] - bbox_a[3]
    return 0 <= salto <= interlineado_max_pt


def extraer_pdf(ruta: str, doc_id: str) -> list[Bloque]:
    """Un bloque = un PARRAFO visual: lineas consecutivas del mismo
    `block` de page.get_text('dict') que comparten firma tipografica
    (fuente/tamano/negrita/cursiva), indentacion x0 compatible, y salto
    vertical de interlineado normal se fusionan en un solo Bloque.

    Motivo del cambio (encontrado en Checkpoint C3a, ver plan y
    forma_ir/segmentacion.py): con un bloque = una linea fisica, un PDF
    de prosa academica en dos columnas fragmentaba cada parrafo en
    decenas de "bloques" de firma casi identica; SEQUITUR agrupaba esa
    repeticion de firma en reglas de gramatica que la Fase C2 marcaba
    como "candidatas a limite" solo por tener diversidad lexica alta
    (trivial en cualquier texto corrido), inundando C3a de falsos
    positivos (27% de lineas marcadas en un PDF real de prosa). Fusionar
    en parrafos ataca la causa raiz: un titulo real de seccion sigue
    siendo su propio bloque corto (nunca se fusiona con el parrafo de
    cuerpo que le sigue, porque difiere en tamano/negrita), mientras que
    el cuerpo de un parrafo pasa a ser UN bloque con UNA firma, no
    docenas de lineas con firmas casi-repetidas."""
    documento = fitz.open(ruta)
    bloques: list[Bloque] = []
    seq = 0
    try:
        for num_pagina, pagina in enumerate(documento, start=1):
            d = pagina.get_text("dict")
            bloque_anterior_y1 = None
            for block in d["blocks"]:
                if "lines" not in block:
                    continue  # bloque de imagen, sin texto

                lineas_utiles = [ln for ln in block["lines"] if "".join(s["text"] for s in ln["spans"]).strip()]
                if not lineas_utiles:
                    continue

                # Agrupa lineas consecutivas del mismo block en parrafos
                grupos: list[list[dict]] = []
                for linea in lineas_utiles:
                    if grupos and _mismo_parrafo(grupos[-1][-1], linea):
                        grupos[-1].append(linea)
                    else:
                        grupos.append([linea])

                for grupo in grupos:
                    texto_parrafo = " ".join(
                        "".join(s["text"] for s in ln["spans"]).strip() for ln in grupo
                    ).strip()
                    if not texto_parrafo:
                        continue

                    primera_linea, ultima_linea = grupo[0], grupo[-1]
                    primer_span = primera_linea["spans"][0]
                    x0 = min(ln["bbox"][0] for ln in grupo)
                    y0 = primera_linea["bbox"][1]
                    x1 = max(ln["bbox"][2] for ln in grupo)
                    y1 = ultima_linea["bbox"][3]
                    bbox = tuple(round(v, 2) for v in (x0, y0, x1, y1))
                    font_size = round(primer_span["size"], 2)
                    bold = _es_negrita(primer_span["flags"], primer_span["font"])
                    italic = _es_italica(primer_span["flags"], primer_span["font"])

                    espacio_vertical = None
                    if bloque_anterior_y1 is not None:
                        espacio_vertical = round(bbox[1] - bloque_anterior_y1, 2)
                    bloque_anterior_y1 = bbox[3]

                    bloques.append(Bloque(
                        bloque_id=f"{doc_id}#p{num_pagina}#{seq}",
                        doc_id=doc_id,
                        seq=seq,
                        texto=texto_parrafo,
                        pagina=num_pagina,
                        diapositiva=None,
                        bbox=bbox,
                        font_size=font_size,
                        bold=bold,
                        italic=italic,
                        indentacion_pt=bbox[0],
                        espacio_vertical_antes=espacio_vertical,
                        formato_fuente="pdf",
                    ))
                    seq += 1
    finally:
        documento.close()
    return bloques


def extraer_pdf_ocr(ruta_json: str, doc_id: str) -> list[Bloque]:
    """Documentos con texto obtenido por OCR (JSON pre-generado, ver
    corpus/_ocr/*.json en produccion). Sin bbox por caja de texto -- se
    degrada a un bloque por pagina, simplificacion explicita (ver plan,
    Fase A). Solo las pistas derivables del texto mismo sobreviven."""
    import json
    with open(ruta_json, encoding="utf-8") as f:
        data = json.load(f)
    bloques: list[Bloque] = []
    for seq, pag in enumerate(data.get("paginas", [])):
        texto = (pag.get("text") or "").strip()
        if not texto:
            continue
        bloques.append(Bloque(
            bloque_id=f"{doc_id}#pocr{pag.get('page', seq+1)}#{seq}",
            doc_id=doc_id,
            seq=seq,
            texto=texto,
            pagina=pag.get("page", seq + 1),
            diapositiva=None,
            bbox=None,
            font_size=None,
            bold=None,
            italic=None,
            indentacion_pt=None,
            espacio_vertical_antes=None,
            formato_fuente="pdf_ocr",
        ))
    return bloques


def extraer_pptx(ruta: str, doc_id: str) -> list[Bloque]:
    """Un bloque = una caja de texto (shape) completa, no cada parrafo
    suelto dentro de ella -- las diapositivas de este corpus suelen tener
    cajas de texto cortas (titulo, lista de vinetas), y separar por
    parrafo perderia la cohesion visual de "una idea por caja"."""
    prs = Presentation(ruta)
    bloques: list[Bloque] = []
    seq = 0
    for num_diapo, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            texto = shape.text_frame.text.strip()
            if not texto:
                continue

            font_size = None
            bold = None
            for parrafo in shape.text_frame.paragraphs:
                for run in parrafo.runs:
                    if run.font.size is not None:
                        font_size = run.font.size.pt
                    if run.font.bold is not None:
                        bold = run.font.bold
                    if font_size is not None:
                        break
                if font_size is not None:
                    break

            try:
                left = Emu(shape.left).pt if shape.left is not None else None
                top = Emu(shape.top).pt if shape.top is not None else None
                width = Emu(shape.width).pt if shape.width is not None else None
                height = Emu(shape.height).pt if shape.height is not None else None
            except (TypeError, ValueError):
                left = top = width = height = None

            bbox = (left, top, left + width, top + height) if None not in (left, top, width, height) else None

            bloques.append(Bloque(
                bloque_id=f"{doc_id}#d{num_diapo}#{seq}",
                doc_id=doc_id,
                seq=seq,
                texto=texto,
                pagina=None,
                diapositiva=num_diapo,
                bbox=bbox,
                font_size=font_size,
                bold=bold,
                italic=None,  # python-pptx no expone italic de forma tan directa como bold; se omite en 1ra pasada
                indentacion_pt=left,
                espacio_vertical_antes=None,  # requeriria trackear el shape anterior en la misma diapositiva; se deja None en 1ra pasada
                formato_fuente="pptx",
            ))
            seq += 1
    return bloques


def extraer_docx(ruta: str, doc_id: str) -> list[Bloque]:
    """Un bloque = un parrafo."""
    d = docx.Document(ruta)
    bloques: list[Bloque] = []
    seq = 0
    for parrafo in d.paragraphs:
        texto = parrafo.text.strip()
        if not texto:
            continue

        indent = None
        if parrafo.paragraph_format.left_indent is not None:
            indent = parrafo.paragraph_format.left_indent.pt

        font_size = None
        bold = None
        italic = None
        for run in parrafo.runs:
            if run.text.strip():
                if run.font.size is not None:
                    font_size = run.font.size.pt
                bold = run.font.bold
                italic = run.font.italic
                break

        bloques.append(Bloque(
            bloque_id=f"{doc_id}#w{seq}",
            doc_id=doc_id,
            seq=seq,
            texto=texto,
            pagina=None,
            diapositiva=None,
            bbox=None,
            font_size=font_size,
            bold=bold,
            italic=italic,
            indentacion_pt=indent,
            espacio_vertical_antes=None,
            formato_fuente="docx",
        ))
        seq += 1
    return bloques
